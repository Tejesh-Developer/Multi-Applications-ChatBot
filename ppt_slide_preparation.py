import os
import streamlit as st
from groq import Groq
from dotenv import load_dotenv
from pptx import Presentation
from pptx.util import Inches

load_dotenv()

def generate_ppt(title, content):
    prs = Presentation()

    slides = content.split("###")

    for slide_text in slides:
        if slide_text.strip() == "":
            continue

        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)

        lines = slide_text.strip().split("\n")
        slide.shapes.title.text = lines[0]

        body = slide.placeholders[1]
        body.text = "\n".join(lines[1:])

    file_name = "generated_presentation.pptx"
    prs.save(file_name)
    return file_name


def run():

    st.markdown("""
        <style>
        .ppt-box {
            background-color: #161B22;
            padding: 20px;
            border-radius: 15px;
            margin-top: 20px;
        }
        </style>
    """, unsafe_allow_html=True)
    # 🔥 HERO BANNER
    st.markdown("""
    <style>
    .hero-banner {
        background: linear-gradient(135deg, #0f2027, #203a43, #2c5364);
        padding: 40px;
        border-radius: 25px;
        margin-bottom: 30px;
        box-shadow: 0 10px 30px rgba(0,0,0,0.4);
    }

    .hero-title {
        font-size: 42px;
        font-weight: 800;
        color: white;
        margin-bottom: 10px;
    }

    .hero-subtitle {
        font-size: 18px;
        color: #cbd5e1;
    }
    </style>

    <div class="hero-banner">
        <div class="hero-title">📊 PPT Slide Preparation</div>
            <div class="hero-subtitle">
                Generate professional PowerPoint slides using AI.
            </div>
        </div>
    """, unsafe_allow_html=True)

    topic = st.text_input("Enter your presentation topic:")

    slide_count = st.slider("Number of Slides", 3, 10, 5)

    if st.button("🚀 Generate Slides") and topic:

        with st.spinner("Generating presentation..."):

            prompt = f"""
            Create a {slide_count}-slide presentation on {topic}.
            Format each slide like:
            ### Slide Title
            Bullet point 1
            Bullet point 2
            Bullet point 3
            """

            client = Groq(api_key=os.getenv("GROQ_API_KEY"))

            chat_completion = client.chat.completions.create(
                messages=[{"role": "user", "content": prompt}],
                model="llama-3.3-70b-versatile",
            )

            ai_text = chat_completion.choices[0].message.content

            ppt_file = generate_ppt(topic, ai_text)

            st.markdown('<div class="ppt-box">', unsafe_allow_html=True)
            st.success("Presentation Generated Successfully!")
            st.download_button(
                "⬇ Download PPT",
                data=open(ppt_file, "rb"),
                file_name="AI_Presentation.pptx"
            )
            st.markdown('</div>', unsafe_allow_html=True)

